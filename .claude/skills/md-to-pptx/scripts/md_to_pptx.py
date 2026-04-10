#!/usr/bin/env python3
"""
md_to_pptx.py — Markdown → PPTX 转换

将预处理后的 Markdown（含 pipe tables 和图片引用）转换为结构化 PowerPoint 文件。
设计风格参考 baoyu-slide-deck（corporate 风格：干净布局、16:9、专业排版）。

用法:
    python md_to_pptx.py input.md -o output.pptx --resource-path=./images_dir
"""

from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Style system ──────────────────────────────────────────────────

@dataclass
class StyleConfig:
    """Complete visual style definition for PPTX generation."""
    name: str
    title_bg: tuple  # RGB tuple (r, g, b) 0-255
    accent: tuple
    text_color: tuple
    bg_color: tuple
    light_bg: tuple
    table_header: tuple
    table_alt: tuple
    font_title: str = "Microsoft YaHei"
    font_body: str = "Microsoft YaHei"

    def rgb(self, attr: str) -> RGBColor:
        t = getattr(self, attr)
        return RGBColor(t[0], t[1], t[2])


def _hex(h: str) -> tuple:
    """Convert '#RRGGBB' to (r, g, b) tuple."""
    h = h.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


STYLES: dict[str, StyleConfig] = {
    "corporate": StyleConfig(
        name="corporate",
        title_bg=_hex("#1E3A5F"), accent=_hex("#2B6CB0"),
        text_color=_hex("#2D2D2D"), bg_color=_hex("#FFFFFF"),
        light_bg=_hex("#F3F4F6"),
        table_header=_hex("#1E3A5F"), table_alt=_hex("#EBEFF5"),
    ),
    "government": StyleConfig(
        name="government",
        title_bg=_hex("#003366"), accent=_hex("#CC0000"),
        text_color=_hex("#333333"), bg_color=_hex("#F5F5F5"),
        light_bg=_hex("#E8E8E8"),
        table_header=_hex("#003366"), table_alt=_hex("#E6EBF0"),
        font_title="SimHei", font_body="SimSun",
    ),
    "education": StyleConfig(
        name="education",
        title_bg=_hex("#2E5090"), accent=_hex("#2FBF71"),
        text_color=_hex("#2C3E50"), bg_color=_hex("#FAFBFC"),
        light_bg=_hex("#EEF2F7"),
        table_header=_hex("#2E5090"), table_alt=_hex("#E8F5E9"),
    ),
    "minimal": StyleConfig(
        name="minimal",
        title_bg=_hex("#555555"), accent=_hex("#0099FF"),
        text_color=_hex("#222222"), bg_color=_hex("#FFFFFF"),
        light_bg=_hex("#F0F0F0"),
        table_header=_hex("#555555"), table_alt=_hex("#F5F5F5"),
    ),
    "technical": StyleConfig(
        name="technical",
        title_bg=_hex("#1A1A1A"), accent=_hex("#0066CC"),
        text_color=_hex("#333333"), bg_color=_hex("#F0F0F0"),
        light_bg=_hex("#E0E0E0"),
        table_header=_hex("#1A1A1A"), table_alt=_hex("#E8E8E8"),
    ),
    "warm": StyleConfig(
        name="warm",
        title_bg=_hex("#B8764F"), accent=_hex("#E67E22"),
        text_color=_hex("#3D3D3D"), bg_color=_hex("#FFFAF5"),
        light_bg=_hex("#F4E4D7"),
        table_header=_hex("#B8764F"), table_alt=_hex("#FFF0E6"),
    ),
}

DEFAULT_STYLE = "corporate"


def load_style(name: str = DEFAULT_STYLE, style_file: str | None = None) -> StyleConfig:
    """Load a style by preset name or from a JSON file."""
    if style_file:
        import json
        data = json.loads(Path(style_file).read_text(encoding="utf-8"))
        return StyleConfig(
            name=data.get("name", "custom"),
            title_bg=_hex(data["title_bg"]), accent=_hex(data["accent"]),
            text_color=_hex(data.get("text_color", "#2D2D2D")),
            bg_color=_hex(data.get("bg_color", "#FFFFFF")),
            light_bg=_hex(data.get("light_bg", "#F3F4F6")),
            table_header=_hex(data.get("table_header", data["title_bg"])),
            table_alt=_hex(data.get("table_alt", "#EBEFF5")),
            font_title=data.get("font_title", "Microsoft YaHei"),
            font_body=data.get("font_body", "Microsoft YaHei"),
        )
    if name not in STYLES:
        print(f"[WARN] 未知风格 '{name}'，使用 {DEFAULT_STYLE}")
        name = DEFAULT_STYLE
    return STYLES[name]


# ── Layout constants (shared across styles) ───────────────────────

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MARGIN_LEFT = Inches(0.8)
MARGIN_TOP = Inches(0.6)
MARGIN_RIGHT = Inches(0.8)
MARGIN_BOTTOM = Inches(0.5)

TITLE_BAR_HEIGHT = Inches(1.0)
BODY_TOP = MARGIN_TOP + TITLE_BAR_HEIGHT + Inches(0.2)
BODY_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
BODY_HEIGHT = SLIDE_HEIGHT - BODY_TOP - MARGIN_BOTTOM

FONT_SIZE_COVER_TITLE = Pt(36)
FONT_SIZE_COVER_SUBTITLE = Pt(18)
FONT_SIZE_SECTION_TITLE = Pt(32)
FONT_SIZE_SLIDE_TITLE = Pt(24)
FONT_SIZE_BODY = Pt(16)
FONT_SIZE_TABLE = Pt(12)
FONT_SIZE_TABLE_HEADER = Pt(13)

MAX_BODY_LINES = 12
MAX_BODY_CHARS = 400
MAX_TABLE_ROWS = 10


# ── Data structures ────────────────────────────────────────────────

@dataclass
class ContentBlock:
    kind: str  # "paragraph", "table", "image", "image_group", "list"
    text: str = ""
    table_data: list[list[str]] = field(default_factory=list)
    image_path: str = ""
    image_paths: list[str] = field(default_factory=list)
    image_layout: str = "horizontal"  # "horizontal", "vertical"


@dataclass
class Section:
    level: int
    title: str
    blocks: list[ContentBlock] = field(default_factory=list)


@dataclass
class SlideData:
    layout: str  # "cover", "section", "content", "table", "image"
    title: str = ""
    subtitle: str = ""
    body_text: str = ""
    table_data: list[list[str]] = field(default_factory=list)
    image_path: str = ""
    image_paths: list[str] = field(default_factory=list)
    image_layout: str = "horizontal"  # "horizontal", "vertical"


# ── Markdown parsing ───────────────────────────────────────────────

def normalize_headings(md_text: str) -> str:
    pattern = re.compile(r"^(#{1,6})([^\s#].*)$", re.MULTILINE)
    return pattern.sub(r"\1 \2", md_text)


def _is_image_line(line: str) -> re.Match | None:
    """Match image lines: ![alt](path) or ![图N] shorthand."""
    return (re.match(r"^!\[([^\]]*)\]\(([^)]+)\)", line) or
            re.match(r"^!\[([^\]]+)\]\s*$", line))


def _is_hr_line(line: str) -> bool:
    """Match horizontal rule: --- (used as vertical layout separator)."""
    return bool(re.match(r"^-{3,}\s*$", line.strip()))


def parse_markdown(md_text: str) -> list[Section]:
    """Parse markdown into a list of Section objects.

    Image handling:
      - ![alt](path)  — standard markdown image with file path
      - ![图1]         — shorthand, resolved to numbered file in resource-path
      - Consecutive image lines → grouped into one image_group block
      - --- between images → vertical layout (default: horizontal)
    """
    md_text = normalize_headings(md_text)
    lines = md_text.split("\n")
    sections: list[Section] = []
    current: Section | None = None
    heading_re = re.compile(r"^(#{1,6})\s+(.+)$")
    table_re = re.compile(r"^\|.+\|$")
    separator_re = re.compile(r"^\|[\s\-:|]+\|$")
    list_re = re.compile(r"^(\s*)([-*]|\d+\.)\s+(.+)$")

    i = 0
    while i < len(lines):
        line = lines[i]

        # Heading
        m = heading_re.match(line)
        if m:
            level = len(m.group(1))
            title = m.group(2).strip()
            current = Section(level=level, title=title)
            sections.append(current)
            i += 1
            continue

        # Ensure a default section for content before any heading
        if current is None:
            current = Section(level=0, title="")
            sections.append(current)

        # Table block (pipe tables)
        if table_re.match(line):
            table_lines = []
            while i < len(lines) and table_re.match(lines[i]):
                if not separator_re.match(lines[i]):
                    row = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                    table_lines.append(row)
                i += 1
            if table_lines:
                current.blocks.append(ContentBlock(kind="table", table_data=table_lines))
            continue

        # Image or image group
        img_m = _is_image_line(line)
        if img_m:
            image_refs: list[str] = []
            layout = "horizontal"

            # Collect consecutive image lines (with optional --- separators)
            while i < len(lines):
                im = _is_image_line(lines[i])
                if im:
                    # Extract path: group(2) for ![alt](path), group(1) for ![图N]
                    path = im.group(2) if im.lastindex and im.lastindex >= 2 else im.group(1)
                    image_refs.append(path)
                    i += 1
                elif _is_hr_line(lines[i]):
                    layout = "vertical"
                    i += 1
                else:
                    break

            if len(image_refs) == 1:
                current.blocks.append(ContentBlock(
                    kind="image", image_path=image_refs[0]))
            else:
                current.blocks.append(ContentBlock(
                    kind="image_group",
                    image_paths=image_refs,
                    image_layout=layout))
            continue

        # List items — accumulate consecutive list lines
        m = list_re.match(line)
        if m:
            list_text_lines = []
            while i < len(lines) and list_re.match(lines[i]):
                lm = list_re.match(lines[i])
                list_text_lines.append(lm.group(3).strip())
                i += 1
            current.blocks.append(ContentBlock(kind="list", text="\n".join(list_text_lines)))
            continue

        # Paragraph — accumulate non-empty lines
        if line.strip():
            para_lines = []
            while i < len(lines) and lines[i].strip() and not heading_re.match(lines[i]) \
                    and not table_re.match(lines[i]) and not _is_image_line(lines[i]) \
                    and not list_re.match(lines[i]):
                para_lines.append(lines[i].strip())
                i += 1
            current.blocks.append(ContentBlock(kind="paragraph", text="\n".join(para_lines)))
            continue

        i += 1  # skip empty lines

    return sections


# ── Slide planning ─────────────────────────────────────────────────

def _estimate_text_lines(text: str) -> int:
    """Estimate how many slide lines a text block will need."""
    lines = 0
    for line in text.split("\n"):
        # ~50 CJK chars or ~80 Latin chars per line
        char_count = sum(2 if ord(c) > 0x2E7F else 1 for c in line)
        lines += max(1, (char_count + 79) // 80)
    return lines


def _split_text(text: str, max_chars: int = MAX_BODY_CHARS, max_lines: int = MAX_BODY_LINES) -> list[str]:
    """Split long text into chunks that fit on a slide."""
    paragraphs = text.split("\n")
    chunks: list[str] = []
    current_chunk: list[str] = []
    current_lines = 0
    current_chars = 0

    for para in paragraphs:
        para_lines = _estimate_text_lines(para)
        para_chars = len(para)

        if current_chunk and (current_lines + para_lines > max_lines or current_chars + para_chars > max_chars):
            chunks.append("\n".join(current_chunk))
            current_chunk = []
            current_lines = 0
            current_chars = 0

        current_chunk.append(para)
        current_lines += para_lines
        current_chars += para_chars

    if current_chunk:
        chunks.append("\n".join(current_chunk))

    return chunks if chunks else [""]


def _split_table(table_data: list[list[str]], max_rows: int = MAX_TABLE_ROWS) -> list[list[list[str]]]:
    """Split a large table into multiple chunks, repeating header."""
    if len(table_data) <= max_rows:
        return [table_data]

    header = table_data[0:1]
    body = table_data[1:]
    chunks = []
    for i in range(0, len(body), max_rows - 1):
        chunk = header + body[i:i + max_rows - 1]
        chunks.append(chunk)
    return chunks


def plan_slides(sections: list[Section]) -> list[SlideData]:
    """Convert sections into a list of SlideData for PPTX generation."""
    slides: list[SlideData] = []
    first_h1 = True

    for section in sections:
        if section.level == 0:
            # Content before any heading — treat as content slides
            for block in section.blocks:
                slides.extend(_block_to_slides(block, ""))
            continue

        if section.level == 1:
            if first_h1:
                slides.append(SlideData(layout="cover", title=section.title))
                first_h1 = False
            else:
                slides.append(SlideData(layout="section", title=section.title))

            # H1 content blocks
            for block in section.blocks:
                slides.extend(_block_to_slides(block, section.title))
        else:
            # H2+ — content slides with the heading as slide title
            slide_title = section.title

            if not section.blocks:
                slides.append(SlideData(layout="content", title=slide_title))
                continue

            for block in section.blocks:
                slides.extend(_block_to_slides(block, slide_title))

    return slides


def _block_to_slides(block: ContentBlock, title: str) -> list[SlideData]:
    """Convert a content block into one or more SlideData."""
    result: list[SlideData] = []

    if block.kind == "table":
        chunks = _split_table(block.table_data)
        for idx, chunk in enumerate(chunks):
            t = title if len(chunks) == 1 else f"{title} (续{idx + 1})" if idx > 0 else title
            result.append(SlideData(layout="table", title=t, table_data=chunk))

    elif block.kind == "image":
        result.append(SlideData(layout="image", title=title, image_path=block.image_path))

    elif block.kind == "image_group":
        result.append(SlideData(
            layout="image", title=title,
            image_paths=block.image_paths,
            image_layout=block.image_layout))

    elif block.kind in ("paragraph", "list"):
        text = block.text
        chunks = _split_text(text)
        for idx, chunk in enumerate(chunks):
            t = title if len(chunks) == 1 else f"{title} (续{idx + 1})" if idx > 0 else title
            result.append(SlideData(layout="content", title=t, body_text=chunk))

    return result


# ── PPTX generation ────────────────────────────────────────────────

def _set_font(run, style: StyleConfig, size=None, bold=False, color_attr="text_color", color_override=None):
    """Apply font settings to a text run."""
    run.font.name = style.font_body
    run.font.size = size or FONT_SIZE_BODY
    run.font.bold = bold
    run.font.color.rgb = color_override or style.rgb(color_attr)


def _add_title_bar(slide, prs, title: str, style: StyleConfig):
    """Add a colored title bar at the top of a slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=0, top=0,
        width=prs.slide_width,
        height=TITLE_BAR_HEIGHT,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = style.rgb("title_bg")
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(
        left=MARGIN_LEFT,
        top=Inches(0.15),
        width=BODY_WIDTH,
        height=TITLE_BAR_HEIGHT - Inches(0.3),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = style.font_title
    run.font.size = FONT_SIZE_SLIDE_TITLE
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _add_cover_slide(prs, slide_data: SlideData, style: StyleConfig):
    """Generate cover slide with centered title."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left=0, top=0,
        width=prs.slide_width, height=prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = style.rgb("title_bg")
    bg.line.fill.background()

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left=Inches(4.0), top=Inches(3.4),
        width=Inches(5.333), height=Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = style.rgb("accent")
    line.line.fill.background()

    txBox = slide.shapes.add_textbox(
        left=Inches(1.5), top=Inches(1.8),
        width=Inches(10.333), height=Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = slide_data.title
    run.font.name = style.font_title
    run.font.size = FONT_SIZE_COVER_TITLE
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if slide_data.subtitle:
        txBox2 = slide.shapes.add_textbox(
            left=Inches(2.0), top=Inches(3.8),
            width=Inches(9.333), height=Inches(1.0))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = slide_data.subtitle
        _set_font(run2, style, size=FONT_SIZE_COVER_SUBTITLE, color_attr="light_bg")


def _add_section_slide(prs, slide_data: SlideData, style: StyleConfig):
    """Generate section divider slide."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left=0, top=Inches(2.2),
        width=prs.slide_width, height=Inches(3.0))
    band.fill.solid()
    band.fill.fore_color.rgb = style.rgb("title_bg")
    band.line.fill.background()

    txBox = slide.shapes.add_textbox(
        left=Inches(1.5), top=Inches(2.8),
        width=Inches(10.333), height=Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = slide_data.title
    run.font.name = style.font_title
    run.font.size = FONT_SIZE_SECTION_TITLE
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _add_content_slide(prs, slide_data: SlideData, style: StyleConfig):
    """Generate content slide with title bar and body text."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    _add_title_bar(slide, prs, slide_data.title, style)

    txBox = slide.shapes.add_textbox(
        left=MARGIN_LEFT, top=BODY_TOP, width=BODY_WIDTH, height=BODY_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = slide_data.body_text.split("\n")
    for idx, line_text in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line_text
        _set_font(run, style)


def _add_table_slide(prs, slide_data: SlideData, style: StyleConfig):
    """Generate table slide with title bar and native table."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    _add_title_bar(slide, prs, slide_data.title, style)

    if not slide_data.table_data:
        return

    rows = len(slide_data.table_data)
    cols = max(len(r) for r in slide_data.table_data)

    table = slide.shapes.add_table(
        rows, cols, left=MARGIN_LEFT, top=BODY_TOP,
        width=BODY_WIDTH, height=BODY_HEIGHT).table

    col_width = int(BODY_WIDTH / cols)
    for c in range(cols):
        table.columns[c].width = col_width

    for r_idx, row in enumerate(slide_data.table_data):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            cell_text = row[c_idx] if c_idx < len(row) else ""
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = cell_text

            if r_idx == 0:
                _set_font(run, style, size=FONT_SIZE_TABLE_HEADER, bold=True,
                          color_override=RGBColor(0xFF, 0xFF, 0xFF))
                cell.fill.solid()
                cell.fill.fore_color.rgb = style.rgb("table_header")
            else:
                _set_font(run, style, size=FONT_SIZE_TABLE)
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = style.rgb("table_alt")

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


IMAGE_EXTENSIONS = (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff")
IMAGE_GAP = Inches(0.15)


def _resolve_image(ref: str, resource_path: str) -> Path | None:
    """Resolve an image reference to a file path.

    Supports:
      - Direct path: "images/firewall.png"
      - Shorthand:   "图1" → searches for 01.*, 1.*, 图1.* in resource_path
    """
    # Try direct path first
    for base in (Path(resource_path), Path(".")):
        candidate = base / ref
        if candidate.exists():
            return candidate

    # Shorthand: extract number from "图N" or "图N-M" or just "N"
    num_match = re.search(r"(\d[\d\-]*)", ref)
    if not num_match:
        return None

    num_str = num_match.group(1)  # e.g. "1", "3-2"
    rp = Path(resource_path)

    # Search patterns: 01.png, 1.png, 图1.png, etc.
    candidates = [
        num_str.zfill(2),           # "01"
        num_str,                     # "1"
        f"图{num_str}",             # "图1"
        f"fig{num_str}",            # "fig1"
        f"figure{num_str}",         # "figure1"
    ]

    if rp.is_dir():
        for f in rp.iterdir():
            if not f.is_file():
                continue
            stem = f.stem.lower()
            for c in candidates:
                if stem == c.lower() and f.suffix.lower() in IMAGE_EXTENSIONS:
                    return f

    return None


def _place_single_image(slide, img_path: Path, left, top, available_w, available_h):
    """Place one image within the given rectangle, centered and aspect-preserved."""
    from PIL import Image as PILImage
    with PILImage.open(img_path) as img:
        img_w, img_h = img.size

    scale = min(available_w / img_w, available_h / img_h)
    display_w = int(img_w * scale)
    display_h = int(img_h * scale)

    cx = left + (available_w - display_w) // 2
    cy = top + (available_h - display_h) // 2

    slide.shapes.add_picture(str(img_path), left=cx, top=cy,
                             width=display_w, height=display_h)


def _add_image_slide(prs, slide_data: SlideData, resource_path: str, style: StyleConfig):
    """Generate image slide with title bar and one or more images."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    _add_title_bar(slide, prs, slide_data.title, style)

    # Collect image paths to render
    refs: list[str] = []
    if slide_data.image_paths:
        refs = slide_data.image_paths
    elif slide_data.image_path:
        refs = [slide_data.image_path]

    resolved: list[Path | None] = [_resolve_image(r, resource_path) for r in refs]

    # Filter to found images, track missing
    missing = [refs[i] for i, p in enumerate(resolved) if p is None]
    found = [(refs[i], p) for i, p in enumerate(resolved) if p is not None]

    if not found:
        # All missing — show placeholder
        txBox = slide.shapes.add_textbox(
            left=MARGIN_LEFT, top=BODY_TOP,
            width=BODY_WIDTH, height=BODY_HEIGHT)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"[图片未找到: {', '.join(refs)}]"
        _set_font(run, style, size=FONT_SIZE_BODY, color_attr="accent")
        return

    n = len(found)
    layout = slide_data.image_layout
    avail_w = BODY_WIDTH
    avail_h = BODY_HEIGHT

    if n == 1:
        _place_single_image(slide, found[0][1], MARGIN_LEFT, BODY_TOP, avail_w, avail_h)

    elif n == 2:
        if layout == "vertical":
            # Top-bottom
            cell_h = (avail_h - IMAGE_GAP) // 2
            _place_single_image(slide, found[0][1],
                                MARGIN_LEFT, BODY_TOP, avail_w, cell_h)
            _place_single_image(slide, found[1][1],
                                MARGIN_LEFT, BODY_TOP + cell_h + IMAGE_GAP, avail_w, cell_h)
        else:
            # Left-right
            cell_w = (avail_w - IMAGE_GAP) // 2
            _place_single_image(slide, found[0][1],
                                MARGIN_LEFT, BODY_TOP, cell_w, avail_h)
            _place_single_image(slide, found[1][1],
                                MARGIN_LEFT + cell_w + IMAGE_GAP, BODY_TOP, cell_w, avail_h)

    elif n == 3:
        if layout == "vertical":
            # Left 1, right 2 (stacked)
            left_w = avail_w // 2 - IMAGE_GAP // 2
            right_w = avail_w - left_w - IMAGE_GAP
            cell_h = (avail_h - IMAGE_GAP) // 2
            _place_single_image(slide, found[0][1],
                                MARGIN_LEFT, BODY_TOP, left_w, avail_h)
            _place_single_image(slide, found[1][1],
                                MARGIN_LEFT + left_w + IMAGE_GAP, BODY_TOP, right_w, cell_h)
            _place_single_image(slide, found[2][1],
                                MARGIN_LEFT + left_w + IMAGE_GAP,
                                BODY_TOP + cell_h + IMAGE_GAP, right_w, cell_h)
        else:
            # Top 1, bottom 2
            top_h = avail_h // 2 - IMAGE_GAP // 2
            bot_h = avail_h - top_h - IMAGE_GAP
            cell_w = (avail_w - IMAGE_GAP) // 2
            _place_single_image(slide, found[0][1],
                                MARGIN_LEFT, BODY_TOP, avail_w, top_h)
            _place_single_image(slide, found[1][1],
                                MARGIN_LEFT, BODY_TOP + top_h + IMAGE_GAP, cell_w, bot_h)
            _place_single_image(slide, found[2][1],
                                MARGIN_LEFT + cell_w + IMAGE_GAP,
                                BODY_TOP + top_h + IMAGE_GAP, cell_w, bot_h)

    else:  # 4+: 2×N grid
        cols = 2
        rows = (n + 1) // 2
        cell_w = (avail_w - IMAGE_GAP * (cols - 1)) // cols
        cell_h = (avail_h - IMAGE_GAP * (rows - 1)) // rows
        for idx, (ref, path) in enumerate(found):
            r, c = divmod(idx, cols)
            x = MARGIN_LEFT + c * (cell_w + IMAGE_GAP)
            y = BODY_TOP + r * (cell_h + IMAGE_GAP)
            _place_single_image(slide, path, x, y, cell_w, cell_h)

    # Show missing image refs as small note
    if missing:
        txBox = slide.shapes.add_textbox(
            left=MARGIN_LEFT, top=SLIDE_HEIGHT - MARGIN_BOTTOM - Inches(0.3),
            width=BODY_WIDTH, height=Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = f"[未找到: {', '.join(missing)}]"
        _set_font(run, style, size=Pt(10), color_attr="accent")


def create_pptx(slides: list[SlideData], output_path: Path,
                 resource_path: str = ".", style: StyleConfig | None = None) -> None:
    """Generate PPTX file from slide data."""
    if style is None:
        style = load_style()

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide_data in slides:
        if slide_data.layout == "cover":
            _add_cover_slide(prs, slide_data, style)
        elif slide_data.layout == "section":
            _add_section_slide(prs, slide_data, style)
        elif slide_data.layout == "content":
            _add_content_slide(prs, slide_data, style)
        elif slide_data.layout == "table":
            _add_table_slide(prs, slide_data, style)
        elif slide_data.layout == "image":
            _add_image_slide(prs, slide_data, resource_path, style)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


# ── CLI ────────────────────────────────────────────────────────────

def parse_args() -> argparse.Namespace:
    style_names = ", ".join(STYLES.keys())
    parser = argparse.ArgumentParser(
        description="Convert Markdown file to PPTX (PowerPoint)."
    )
    parser.add_argument("input_md", type=Path, help="Input markdown file path")
    parser.add_argument(
        "-o", "--output", type=Path, default=None,
        help="Output pptx file path (default: same name as input)",
    )
    parser.add_argument(
        "--resource-path", default=".",
        help="Path to resolve image references (default: .)",
    )
    parser.add_argument(
        "--style", default=DEFAULT_STYLE,
        help=f"Style preset ({style_names}) (default: {DEFAULT_STYLE})",
    )
    parser.add_argument(
        "--style-file", default=None,
        help="Custom style JSON file (from extract_style.py)",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    input_md: Path = args.input_md

    if not input_md.exists():
        print(f"[ERROR] 输入文件不存在: {input_md}")
        return 1

    output_pptx = args.output if args.output else input_md.with_suffix(".pptx")

    try:
        style = load_style(args.style, args.style_file)
        md_text = input_md.read_text(encoding="utf-8")
        sections = parse_markdown(md_text)
        slides = plan_slides(sections)
        create_pptx(slides, output_pptx, args.resource_path, style)
    except Exception as e:
        print(f"[ERROR] {e}")
        return 1

    print(f"[OK] 转换完成: {output_pptx} ({len(slides)} 张幻灯片) [风格: {style.name}]")
    return 0


if __name__ == "__main__":
    sys.exit(main())
