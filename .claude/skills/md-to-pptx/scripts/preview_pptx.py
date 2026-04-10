#!/usr/bin/env python3
"""
preview_pptx.py — 在终端预览 PPTX 内容（无需打开文件）

用法:
    python preview_pptx.py demo.pptx
"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

def preview(pptx_path: str):
    prs = Presentation(pptx_path)
    total = len(prs.slides)
    w_in = prs.slide_width / 914400
    h_in = prs.slide_height / 914400

    print(f"{'═' * 60}")
    print(f"  {Path(pptx_path).name}  |  {total} slides  |  {w_in:.1f}×{h_in:.1f} in")
    print(f"{'═' * 60}")

    for i, slide in enumerate(prs.slides, 1):
        # Detect layout type
        has_table = False
        has_image = False
        full_bg = False
        mid_band = False

        for s in slide.shapes:
            if s.has_table:
                has_table = True
            if s.shape_type == 13:  # picture
                has_image = True
            if hasattr(s, 'left') and hasattr(s, 'width'):
                if s.left == 0 and s.width == prs.slide_width:
                    if s.height == prs.slide_height:
                        full_bg = True
                    elif s.height > Emu(1500000) and s.top > Emu(1000000):
                        mid_band = True

        if full_bg:
            layout = "封面"
        elif mid_band:
            layout = "章节"
        elif has_table:
            layout = "表格"
        elif has_image:
            layout = "图片"
        else:
            layout = "内容"

        print(f"\n{'─' * 60}")
        print(f"  Slide {i:02d}/{total}  [{layout}]")
        print(f"{'─' * 60}")

        # Extract text content
        for s in slide.shapes:
            if s.has_text_frame:
                text = s.text_frame.text.strip()
                if text:
                    # Indent multi-line text
                    for line in text.split('\n'):
                        if line.strip():
                            print(f"  {line.strip()}")

            if s.has_table:
                t = s.table
                col_widths = []
                # Calculate column widths for alignment
                for c in range(len(t.columns)):
                    max_w = 0
                    for r in range(len(t.rows)):
                        cell_text = t.cell(r, c).text.strip()
                        # CJK chars count as 2 width
                        w = sum(2 if ord(ch) > 0x2E7F else 1 for ch in cell_text)
                        max_w = max(max_w, w)
                    col_widths.append(min(max_w, 20))

                for r in range(len(t.rows)):
                    cells = []
                    for c in range(len(t.columns)):
                        cell_text = t.cell(r, c).text.strip()[:18]
                        cells.append(cell_text)
                    row_str = " │ ".join(cells)
                    prefix = "  ▸ " if r == 0 else "    "
                    print(f"{prefix}{row_str}")
                    if r == 0:
                        print(f"    {'─' * len(row_str)}")

    print(f"\n{'═' * 60}")
    print(f"  预览完成: {total} 张幻灯片")
    print(f"{'═' * 60}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python preview_pptx.py <file.pptx>")
        sys.exit(1)
    preview(sys.argv[1])
