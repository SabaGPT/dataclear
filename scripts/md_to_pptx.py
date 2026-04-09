#!/usr/bin/env python3
"""
md_to_pptx.py — Markdown → PPTX 转换

将预处理后的 Markdown（含 pipe tables 和图片引用）转换为结构化 PowerPoint 文件。
设计风格参考 baoyu-slide-deck（corporate 风格：干净布局、16:9、专业排版）。

用法:
    python md_to_pptx.py input.md -o output.pptx --resource-path=./images_dir
"""

from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Style constants (baoyu corporate) ──────────────────────────────

COLOR_TITLE_BG = RGBColor(0x1E, 0x3A, 0x5F)     # 深蓝 navy
COLOR_ACCENT = RGBColor(0x2B, 0x6C, 0xB0)         # 蓝色强调
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK = RGBColor(0x2D, 0x2D, 0x2D)           # 正文深灰
COLOR_LIGHT_BG = RGBColor(0xF3, 0xF4, 0xF6)       # 浅灰背景
COLOR_TABLE_HEADER = RGBColor(0x1E, 0x3A, 0x5F)   # 表头背景
COLOR_TABLE_ALT = RGBColor(0xEB, 0xEF, 0xF5)      # 表格交替行

FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"
FONT_SIZE_COVER_TITLE = Pt(36)
FONT_SIZE_COVER_SUBTITLE = Pt(18)
FONT_SIZE_SECTION_TITLE = Pt(32)
FONT_SIZE_SLIDE_TITLE = Pt(24)
FONT_SIZE_BODY = Pt(16)
FONT_SIZE_TABLE = Pt(12)
FONT_SIZE_TABLE_HEADER = Pt(13)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MARGIN_LEFT = Inches(0.8)
MARGIN_TOP = Inches(0.6)
MARGIN_RIGHT = Inches(0.8)
MARGIN_BOTTOM = Inches(0.5)

TITLE_BAR_HEIGHT = Inches(1.0)
BODY_TOP = MARGIN_TOP + TITLE_BAR_HEIGHT + Inches(0.2)
BODY_WIDTH = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
BODY_HEIGHT = SLIDE_HEIGHT - BODY_TOP - MARGIN_BOTTOM

MAX_BODY_LINES = 12
MAX_BODY_CHARS = 400
MAX_TABLE_ROWS = 10


# ── Data structures ────────────────────────────────────────────────

@dataclass
class ContentBlock:
    kind: str  # "paragraph", "table", "image", "list"
    text: str = ""
    table_data: list[list[str]] = field(default_factory=list)
    image_path: str = ""


@dataclass
class Section:
    level: int
    title: str
    blocks: list[ContentBlock] = field(default_factory=list)


@dataclass
class SlideData:
    layout: str  # "cover", "section", "content", "table", "image"
    title: str = ""
    subtitle: str = ""
    body_text: str = ""
    table_data: list[list[str]] = field(default_factory=list)
    image_path: str = ""


# ── Markdown parsing ───────────────────────────────────────────────

def normalize_headings(md_text: str) -> str:
    pattern = re.compile(r"^(#{1,6})([^\s#].*)$", re.MULTILINE)
    return pattern.sub(r"\1 \2", md_text)


def parse_markdown(md_text: str) -> list[Section]:
    """Parse markdown into a list of Section objects."""
    md_text = normalize_headings(md_text)
    lines = md_text.split("\n")
    sections: list[Section] = []
    current: Section | None = None
    heading_re = re.compile(r"^(#{1,6})\s+(.+)$")
    table_re = re.compile(r"^\|.+\|$")
    separator_re = re.compile(r"^\|[\s\-:|]+\|$")
    image_re = re.compile(r"^!\[([^\]]*)\]\(([^)]+)\)")
    list_re = re.compile(r"^(\s*)([-*]|\d+\.)\s+(.+)$")

    i = 0
    while i < len(lines):
        line = lines[i]

        # Heading
        m = heading_re.match(line)
        if m:
            level = len(m.group(1))
            title = m.group(2).strip()
            current = Section(level=level, title=title)
            sections.append(current)
            i += 1
            continue

        # Ensure a default section for content before any heading
        if current is None:
            current = Section(level=0, title="")
            sections.append(current)

        # Table block (pipe tables)
        if table_re.match(line):
            table_lines = []
            while i < len(lines) and table_re.match(lines[i]):
                if not separator_re.match(lines[i]):
                    row = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                    table_lines.append(row)
                i += 1
            if table_lines:
                current.blocks.append(ContentBlock(kind="table", table_data=table_lines))
            continue

        # Image
        m = image_re.match(line)
        if m:
            current.blocks.append(ContentBlock(kind="image", image_path=m.group(2)))
            i += 1
            continue

        # List items — accumulate consecutive list lines
        m = list_re.match(line)
        if m:
            list_text_lines = []
            while i < len(lines) and list_re.match(lines[i]):
                lm = list_re.match(lines[i])
                list_text_lines.append(lm.group(3).strip())
                i += 1
            current.blocks.append(ContentBlock(kind="list", text="\n".join(list_text_lines)))
            continue

        # Paragraph — accumulate non-empty lines
        if line.strip():
            para_lines = []
            while i < len(lines) and lines[i].strip() and not heading_re.match(lines[i]) \
                    and not table_re.match(lines[i]) and not image_re.match(lines[i]) \
                    and not list_re.match(lines[i]):
                para_lines.append(lines[i].strip())
                i += 1
            current.blocks.append(ContentBlock(kind="paragraph", text="\n".join(para_lines)))
            continue

        i += 1  # skip empty lines

    return sections


# ── Slide planning ─────────────────────────────────────────────────

def _estimate_text_lines(text: str) -> int:
    """Estimate how many slide lines a text block will need."""
    lines = 0
    for line in text.split("\n"):
        # ~50 CJK chars or ~80 Latin chars per line
        char_count = sum(2 if ord(c) > 0x2E7F else 1 for c in line)
        lines += max(1, (char_count + 79) // 80)
    return lines


def _split_text(text: str, max_chars: int = MAX_BODY_CHARS, max_lines: int = MAX_BODY_LINES) -> list[str]:
    """Split long text into chunks that fit on a slide."""
    paragraphs = text.split("\n")
    chunks: list[str] = []
    current_chunk: list[str] = []
    current_lines = 0
    current_chars = 0

    for para in paragraphs:
        para_lines = _estimate_text_lines(para)
        para_chars = len(para)

        if current_chunk and (current_lines + para_lines > max_lines or current_chars + para_chars > max_chars):
            chunks.append("\n".join(current_chunk))
            current_chunk = []
            current_lines = 0
            current_chars = 0

        current_chunk.append(para)
        current_lines += para_lines
        current_chars += para_chars

    if current_chunk:
        chunks.append("\n".join(current_chunk))

    return chunks if chunks else [""]


def _split_table(table_data: list[list[str]], max_rows: int = MAX_TABLE_ROWS) -> list[list[list[str]]]:
    """Split a large table into multiple chunks, repeating header."""
    if len(table_data) <= max_rows:
        return [table_data]

    header = table_data[0:1]
    body = table_data[1:]
    chunks = []
    for i in range(0, len(body), max_rows - 1):
        chunk = header + body[i:i + max_rows - 1]
        chunks.append(chunk)
    return chunks


def plan_slides(sections: list[Section]) -> list[SlideData]:
    """Convert sections into a list of SlideData for PPTX generation."""
    slides: list[SlideData] = []
    first_h1 = True

    for section in sections:
        if section.level == 0:
            # Content before any heading — treat as content slides
            for block in section.blocks:
                slides.extend(_block_to_slides(block, ""))
            continue

        if section.level == 1:
            if first_h1:
                slides.append(SlideData(layout="cover", title=section.title))
                first_h1 = False
            else:
                slides.append(SlideData(layout="section", title=section.title))

            # H1 content blocks
            for block in section.blocks:
                slides.extend(_block_to_slides(block, section.title))
        else:
            # H2+ — content slides with the heading as slide title
            slide_title = section.title

            if not section.blocks:
                slides.append(SlideData(layout="content", title=slide_title))
                continue

            for block in section.blocks:
                slides.extend(_block_to_slides(block, slide_title))

    return slides


def _block_to_slides(block: ContentBlock, title: str) -> list[SlideData]:
    """Convert a content block into one or more SlideData."""
    result: list[SlideData] = []

    if block.kind == "table":
        chunks = _split_table(block.table_data)
        for idx, chunk in enumerate(chunks):
            t = title if len(chunks) == 1 else f"{title} (续{idx + 1})" if idx > 0 else title
            result.append(SlideData(layout="table", title=t, table_data=chunk))

    elif block.kind == "image":
        result.append(SlideData(layout="image", title=title, image_path=block.image_path))

    elif block.kind in ("paragraph", "list"):
        text = block.text
        chunks = _split_text(text)
        for idx, chunk in enumerate(chunks):
            t = title if len(chunks) == 1 else f"{title} (续{idx + 1})" if idx > 0 else title
            result.append(SlideData(layout="content", title=t, body_text=chunk))

    return result


# ── PPTX generation ────────────────────────────────────────────────

def _set_font(run, name=FONT_BODY, size=FONT_SIZE_BODY, bold=False, color=COLOR_DARK):
    """Apply font settings to a text run."""
    run.font.name = name
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def _add_title_bar(slide, prs, title: str):
    """Add a colored title bar at the top of a slide."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=0, top=0,
        width=prs.slide_width,
        height=TITLE_BAR_HEIGHT,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_TITLE_BG
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(
        left=MARGIN_LEFT,
        top=Inches(0.15),
        width=BODY_WIDTH,
        height=TITLE_BAR_HEIGHT - Inches(0.3),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    _set_font(run, size=FONT_SIZE_SLIDE_TITLE, bold=True, color=COLOR_WHITE)


def _add_cover_slide(prs, slide_data: SlideData):
    """Generate cover slide with centered title."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Full-slide background color
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=0, top=0,
        width=prs.slide_width,
        height=prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_TITLE_BG
    bg.line.fill.background()

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(4.0),
        top=Inches(3.4),
        width=Inches(5.333),
        height=Inches(0.06),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(
        left=Inches(1.5),
        top=Inches(1.8),
        width=Inches(10.333),
        height=Inches(1.5),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = slide_data.title
    _set_font(run, size=FONT_SIZE_COVER_TITLE, bold=True, color=COLOR_WHITE)

    # Subtitle
    if slide_data.subtitle:
        txBox2 = slide.shapes.add_textbox(
            left=Inches(2.0),
            top=Inches(3.8),
            width=Inches(9.333),
            height=Inches(1.0),
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = slide_data.subtitle
        _set_font(run2, size=FONT_SIZE_COVER_SUBTITLE, color=COLOR_LIGHT_BG)


def _add_section_slide(prs, slide_data: SlideData):
    """Generate section divider slide."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Background band (partial height, centered)
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=0,
        top=Inches(2.2),
        width=prs.slide_width,
        height=Inches(3.0),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = COLOR_TITLE_BG
    band.line.fill.background()

    # Section title
    txBox = slide.shapes.add_textbox(
        left=Inches(1.5),
        top=Inches(2.8),
        width=Inches(10.333),
        height=Inches(1.5),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = slide_data.title
    _set_font(run, size=FONT_SIZE_SECTION_TITLE, bold=True, color=COLOR_WHITE)


def _add_content_slide(prs, slide_data: SlideData):
    """Generate content slide with title bar and body text."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    _add_title_bar(slide, prs, slide_data.title)

    # Body text
    txBox = slide.shapes.add_textbox(
        left=MARGIN_LEFT,
        top=BODY_TOP,
        width=BODY_WIDTH,
        height=BODY_HEIGHT,
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = slide_data.body_text.split("\n")
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        _set_font(run)


def _add_table_slide(prs, slide_data: SlideData):
    """Generate table slide with title bar and native table."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    _add_title_bar(slide, prs, slide_data.title)

    if not slide_data.table_data:
        return

    rows = len(slide_data.table_data)
    cols = max(len(r) for r in slide_data.table_data)

    # Calculate available space
    table_top = BODY_TOP
    table_width = BODY_WIDTH
    table_height = BODY_HEIGHT

    table = slide.shapes.add_table(
        rows, cols,
        left=MARGIN_LEFT,
        top=table_top,
        width=table_width,
        height=table_height,
    ).table

    # Distribute column widths evenly
    col_width = int(table_width / cols)
    for c in range(cols):
        table.columns[c].width = col_width

    for r_idx, row in enumerate(slide_data.table_data):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            cell_text = row[c_idx] if c_idx < len(row) else ""
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = cell_text

            if r_idx == 0:
                # Header row styling
                _set_font(run, size=FONT_SIZE_TABLE_HEADER, bold=True, color=COLOR_WHITE)
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TABLE_HEADER
            else:
                _set_font(run, size=FONT_SIZE_TABLE, color=COLOR_DARK)
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_TABLE_ALT

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_image_slide(prs, slide_data: SlideData, resource_path: str):
    """Generate image slide with title bar and centered image."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    _add_title_bar(slide, prs, slide_data.title)

    # Resolve image path
    img_path = Path(resource_path) / slide_data.image_path
    if not img_path.exists():
        # Try relative to current dir
        img_path = Path(slide_data.image_path)

    if not img_path.exists():
        # Add a text placeholder for missing image
        txBox = slide.shapes.add_textbox(
            left=MARGIN_LEFT, top=BODY_TOP,
            width=BODY_WIDTH, height=BODY_HEIGHT,
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"[图片未找到: {slide_data.image_path}]"
        _set_font(run, size=FONT_SIZE_BODY, color=COLOR_ACCENT)
        return

    # Calculate image size to fit in body area, maintaining aspect ratio
    from PIL import Image
    with Image.open(img_path) as img:
        img_w, img_h = img.size

    available_w = BODY_WIDTH
    available_h = BODY_HEIGHT

    # Scale to fit
    scale_w = available_w / img_w
    scale_h = available_h / img_h
    scale = min(scale_w, scale_h)

    display_w = int(img_w * scale)
    display_h = int(img_h * scale)

    # Center in body area
    left = MARGIN_LEFT + (available_w - display_w) // 2
    top = BODY_TOP + (available_h - display_h) // 2

    slide.shapes.add_picture(
        str(img_path),
        left=left, top=top,
        width=display_w, height=display_h,
    )


def create_pptx(slides: list[SlideData], output_path: Path, resource_path: str = ".") -> None:
    """Generate PPTX file from slide data."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for slide_data in slides:
        if slide_data.layout == "cover":
            _add_cover_slide(prs, slide_data)
        elif slide_data.layout == "section":
            _add_section_slide(prs, slide_data)
        elif slide_data.layout == "content":
            _add_content_slide(prs, slide_data)
        elif slide_data.layout == "table":
            _add_table_slide(prs, slide_data)
        elif slide_data.layout == "image":
            _add_image_slide(prs, slide_data, resource_path)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


# ── CLI ────────────────────────────────────────────────────────────

def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Convert Markdown file to PPTX (PowerPoint)."
    )
    parser.add_argument("input_md", type=Path, help="Input markdown file path")
    parser.add_argument(
        "-o", "--output", type=Path, default=None,
        help="Output pptx file path (default: same name as input)",
    )
    parser.add_argument(
        "--resource-path", default=".",
        help="Path to resolve image references (default: .)",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    input_md: Path = args.input_md

    if not input_md.exists():
        print(f"[ERROR] 输入文件不存在: {input_md}")
        return 1

    output_pptx = args.output if args.output else input_md.with_suffix(".pptx")

    try:
        md_text = input_md.read_text(encoding="utf-8")
        sections = parse_markdown(md_text)
        slides = plan_slides(sections)
        create_pptx(slides, output_pptx, args.resource_path)
    except Exception as e:
        print(f"[ERROR] {e}")
        return 1

    print(f"[OK] 转换完成: {output_pptx} ({len(slides)} 张幻灯片)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
